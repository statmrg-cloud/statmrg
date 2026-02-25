"""
전자책 문서 생성기 (v2 - 버그 수정 + 고급 디자인)
버그 수정:
  - pop from empty list: _draw_paragraph 에서 saveState 제거 (showPage 교차 불가)
  - 빈 페이지: _end_page_if_drawn()로 콘텐츠 있을 때만 showPage
  - 한글 텍스트 넘침: c.stringWidth() 기반 정확한 줄바꿈
개선:
  - 아름다운 프로그래매틱 표지 (외부 이미지 불필요)
  - 고급 목차 레이아웃
  - 챕터 헤더 디자인
PDF, DOCX, PPTX 지원
"""
import os
import re
import sys
import math
import requests
from PIL import Image as PILImage
from reportlab.lib.pagesizes import A4
from reportlab.lib.colors import HexColor, Color
from reportlab.pdfgen import canvas
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

from config import FONT_PATHS, load_config


# ============================================================
# 유틸리티
# ============================================================
def find_font_path(font_name):
    """폰트 파일 경로 검색 (macOS / Windows 크로스 플랫폼)"""
    # 1) 명시적 FONT_PATHS 매핑
    if font_name in FONT_PATHS:
        path = FONT_PATHS[font_name]
        if os.path.exists(path):
            return path

    # 2) 시스템 폰트 디렉토리 검색
    search_dirs = [
        # macOS
        '/System/Library/Fonts/Supplemental/',
        '/System/Library/Fonts/',
        '/Library/Fonts/',
        os.path.expanduser('~/Library/Fonts/'),
        # Windows
        'C:/Windows/Fonts/',
        os.path.expanduser('~/AppData/Local/Microsoft/Windows/Fonts/'),
    ]
    for d in search_dirs:
        if os.path.isdir(d):
            for f in os.listdir(d):
                if font_name.lower() in f.lower() and f.endswith(('.ttf', '.ttc', '.otf')):
                    return os.path.join(d, f)

    # 3) PyInstaller 번들 내 폰트 검색
    if getattr(sys, 'frozen', False):
        bundle_dir = sys._MEIPASS
    else:
        bundle_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    for fname in os.listdir(bundle_dir) if os.path.isdir(bundle_dir) else []:
        if font_name.lower() in fname.lower() and fname.endswith(('.ttf', '.ttc', '.otf')):
            return os.path.join(bundle_dir, fname)

    return None


# Windows 시스템 한국어 폰트 우선순위 (항상 존재)
_WIN_FALLBACK_FONTS = [
    ('MalgunGothic', 'C:/Windows/Fonts/malgun.ttf'),
    ('Batang',       'C:/Windows/Fonts/batang.ttc'),
    ('Gulim',        'C:/Windows/Fonts/gulim.ttc'),
    ('Dotum',        'C:/Windows/Fonts/dotum.ttc'),
]
# macOS 한국어 폰트 우선순위
_MAC_FALLBACK_FONTS = [
    ('AppleGothic',   '/System/Library/Fonts/Supplemental/AppleGothic.ttf'),
    ('AppleMyungjo',  '/System/Library/Fonts/Supplemental/AppleMyungjo.ttf'),
    ('Helvetica',     None),  # PDF 내장 폰트 (영문 fallback)
]

def _register_fallback_font():
    """시스템 기본 한국어 폰트 자동 등록 및 이름 반환"""
    import platform
    candidates = _WIN_FALLBACK_FONTS if platform.system() == 'Windows' else _MAC_FALLBACK_FONTS
    for name, path in candidates:
        if path is None:
            return name  # 내장 폰트
        if os.path.exists(path):
            try:
                pdfmetrics.registerFont(TTFont(name, path))
            except Exception:
                try:
                    pdfmetrics.registerFont(TTFont(name, path, subfontIndex=0))
                except Exception:
                    continue
            return name
    return 'Helvetica'  # 최후 fallback (영문만 지원)


def register_font(font_name):
    """폰트 등록 — 지정 폰트 없으면 시스템 fallback 자동 사용"""
    path = find_font_path(font_name)
    if path:
        try:
            pdfmetrics.registerFont(TTFont(font_name, path))
            return font_name
        except Exception:
            try:
                pdfmetrics.registerFont(TTFont(font_name, path, subfontIndex=0))
                return font_name
            except Exception:
                pass
    # 지정 폰트 없음 → 시스템 기본 폰트로 fallback
    fallback = _register_fallback_font()
    return fallback


def download_image(url, output_dir):
    if not url:
        return None
    try:
        resp = requests.get(url, timeout=30, headers={'User-Agent': 'Mozilla/5.0'})
        if resp.status_code != 200:
            return None
        ext = '.png' if 'png' in resp.headers.get('content-type', '') else '.jpg'
        filename = f"img_{abs(hash(url)) % 100000}{ext}"
        filepath = os.path.join(output_dir, filename)
        with open(filepath, 'wb') as f:
            f.write(resp.content)
        return filepath
    except Exception as e:
        print(f"[PDF] 이미지 다운로드 실패: {e}")
        return None


def _safe_text(text):
    """이모지/미지원 특수문자 → 공백 대체"""
    if not text:
        return ''
    cleaned = []
    for ch in text:
        cp = ord(ch)
        if (cp < 0x10000
                and not (0x1F000 <= cp <= 0x1FFFF)
                and not (0x2600 <= cp <= 0x27BF)
                and not (0xFE00 <= cp <= 0xFE0F)
                and cp not in range(0x200B, 0x200E)):
            cleaned.append(ch)
        else:
            cleaned.append(' ')
    return ''.join(cleaned)


PHASE_COLORS = {
    '문제인식': '#e74c3c',
    '방법발견': '#f39c12',
    '실행':    '#03c75a',
    '확신':    '#3498db',
}

PHASE_BG = {
    '문제인식': '#fff5f5',
    '방법발견': '#fffbf0',
    '실행':    '#f0fff8',
    '확신':    '#f0f8ff',
}


# ============================================================
# PDF 생성기 메인
# ============================================================
class EbookPDFGenerator:
    def __init__(self, config=None):
        self.config = config or load_config()
        self.font_name = register_font(self.config.get('pdf_font', 'AppleGothic'))
        self.font_size = self.config.get('pdf_font_size', 11)
        self.heading_size = self.config.get('pdf_heading_size', 16)
        self.subheading_size = self.config.get('pdf_subheading_size', 13)
        self.line_spacing = self.config.get('pdf_line_spacing', 1.6)
        self.margin_top = self.config.get('pdf_margin_top', 72)
        self.margin_bottom = self.config.get('pdf_margin_bottom', 72)
        self.margin_left = self.config.get('pdf_margin_left', 60)
        self.margin_right = self.config.get('pdf_margin_right', 60)
        self.page_width, self.page_height = A4
        self.content_width = self.page_width - self.margin_left - self.margin_right
        self.output_dir = self.config.get('output_dir', './static/output')
        os.makedirs(self.output_dir, exist_ok=True)
        self._page_drawn = False  # 현재 페이지에 콘텐츠가 있는지 추적
        self.chapter_pages = {}   # 챕터 번호 → 시작 페이지 (2-pass TOC 페이지 번호용)

    # ──────────────────────────────────────────────
    # 페이지 관리
    # ──────────────────────────────────────────────
    def _mark_drawn(self):
        self._page_drawn = True

    def _draw_page_number(self, c):
        page_num = c.getPageNumber()
        if page_num > 1:  # 표지 제외
            c.setFont('Helvetica', 8)
            c.setFillColor(HexColor('#aaaaaa'))
            c.drawCentredString(self.page_width / 2, 22, str(page_num))

    def _new_page(self, c):
        """페이지 번호 그리고 새 페이지 시작 (항상)"""
        self._draw_page_number(c)
        c.showPage()
        self.y = self.page_height - self.margin_top
        self._page_drawn = False

    def _end_page_if_drawn(self, c):
        """콘텐츠가 있을 때만 페이지 종료 → 빈 페이지 방지"""
        has_content = self._page_drawn or (self.y < self.page_height - self.margin_top - 8)
        if has_content:
            self._draw_page_number(c)
            c.showPage()
        self.y = self.page_height - self.margin_top
        self._page_drawn = False

    def _check_page_break(self, c, needed_height=50):
        if self.y - needed_height < self.margin_bottom:
            self._new_page(c)
            return True
        return False

    # ──────────────────────────────────────────────
    # 텍스트 그리기 (핵심 버그 수정 위치)
    # ──────────────────────────────────────────────
    def _wrap_text_by_width(self, c, text, font_size, max_width):
        """실제 폰트 메트릭 기반 줄바꿈 (한글 정확 처리)"""
        if not text:
            return []
        lines = []
        current_line = ''
        for char in text:
            test = current_line + char
            if c.stringWidth(test, self.font_name, font_size) > max_width and current_line:
                # 마지막 공백에서 영어 단어 경계 처리
                last_space = current_line.rfind(' ')
                if last_space > max(0, len(current_line) - 12):
                    lines.append(current_line[:last_space])
                    current_line = current_line[last_space + 1:] + char
                else:
                    lines.append(current_line)
                    current_line = char
            else:
                current_line = test
        if current_line:
            lines.append(current_line)
        return lines or [text]

    def _draw_text(self, c, text, x=None, font_size=None, color='#333333'):
        """한 줄 텍스트 (saveState 없음 - showPage 간섭 없음)"""
        text = _safe_text(text)
        if not text.strip():
            return
        if x is None:
            x = self.margin_left
        if font_size is None:
            font_size = self.font_size
        c.setFont(self.font_name, font_size)
        c.setFillColor(HexColor(color))
        c.drawString(x, self.y, text)
        self.y -= font_size * self.line_spacing
        self._mark_drawn()

    def _draw_text_centered(self, c, text, font_size=None, color='#333333'):
        text = _safe_text(text)
        if not text.strip():
            return
        if font_size is None:
            font_size = self.font_size
        c.setFont(self.font_name, font_size)
        c.setFillColor(HexColor(color))
        c.drawCentredString(self.page_width / 2, self.y, text)
        self.y -= font_size * self.line_spacing
        self._mark_drawn()

    def _draw_paragraph(self, c, text, font_size=None, color='#333333', indent=0):
        """
        자동 줄바꿈 문단 그리기
        ★ 핵심 버그 수정: 외부 saveState/restoreState 제거
          (내부에서 _check_page_break → showPage 가 호출되면
           saveState 스택이 리셋되어 restoreState 시 'pop from empty list' 발생)
        ★ 한글 너비 수정: c.stringWidth() 사용
        """
        text = _safe_text(text)
        if not text.strip():
            return
        if font_size is None:
            font_size = self.font_size

        available_width = self.content_width - indent
        line_height = font_size * self.line_spacing

        for input_line in text.split('\n'):
            stripped = input_line.strip()
            if not stripped:
                self.y -= line_height * 0.5
                continue

            # 실제 폰트 메트릭으로 줄바꿈 (한글 오버플로 방지)
            wrapped = self._wrap_text_by_width(c, stripped, font_size, available_width)
            for wline in wrapped:
                self._check_page_break(c, line_height + 4)
                # ★ 매 줄마다 폰트/색상 재설정 (showPage 이후 초기화 대비)
                c.setFont(self.font_name, font_size)
                c.setFillColor(HexColor(color))
                c.drawString(self.margin_left + indent, self.y, wline)
                self.y -= line_height
                self._mark_drawn()
        # 문단 후 소폭 여백
        self.y -= line_height * 0.15

    # ──────────────────────────────────────────────
    # 메인 생성
    # ──────────────────────────────────────────────
    def _render_pages(self, c, ebook_data, skip_images=False):
        """공통 렌더링 루틴 — skip_images=True이면 챕터 시작 이미지 생략 (1st pass용)"""
        prologue = ebook_data.get('prologue', '')
        epilogue = ebook_data.get('epilogue', '')
        chapters_content = ebook_data.get('chapters_content', [])
        chapter_images   = ebook_data.get('chapter_images', [])

        # 1. 표지
        self._draw_cover(c, ebook_data)
        c.showPage()

        # 2. 목차
        self.y = self.page_height - self.margin_top
        self._page_drawn = False
        self._draw_toc(c, ebook_data)
        self._end_page_if_drawn(c)

        # 3. 프롤로그
        if prologue and prologue.strip():
            self.y = self.page_height - self.margin_top
            self._page_drawn = False
            self._draw_prologue(c, prologue)
            self._end_page_if_drawn(c)

        # 4. 가치 요약
        self.y = self.page_height - self.margin_top
        self._page_drawn = False
        self._draw_analysis_summary(c, ebook_data)
        self._end_page_if_drawn(c)

        # 5. 챕터들
        for i, ch_data in enumerate(chapters_content):
            chapter = ch_data.get('chapter', {})
            content = ch_data.get('content', '')
            img_url = (chapter_images[i] if i < len(chapter_images) else None) if not skip_images else None
            num = i + 1

            self.y = self.page_height - self.margin_top
            self._page_drawn = False
            self.chapter_pages[num] = c.getPageNumber()   # 챕터 시작 페이지 기록
            self._draw_chapter_start(c, chapter, num, img_url)
            self._end_page_if_drawn(c)

            self.y = self.page_height - self.margin_top
            self._page_drawn = False
            self._draw_chapter_content(c, content, num)
            self._end_page_if_drawn(c)

        # 6. 에필로그
        if epilogue and epilogue.strip():
            self.y = self.page_height - self.margin_top
            self._page_drawn = False
            self._draw_epilogue(c, epilogue)
            self._end_page_if_drawn(c)

        # 7. 마케팅 부록
        self.y = self.page_height - self.margin_top
        self._page_drawn = False
        self._draw_marketing_page(c, ebook_data)
        self._end_page_if_drawn(c)

    def generate(self, ebook_data):
        import io as _io
        book_info = ebook_data.get('book_info', {})
        title = book_info.get('book_title', '전자책')
        safe_title = re.sub(r'[^\w가-힣\s-]', '', title)[:50].strip()
        filename = f"{safe_title}.pdf"
        filepath = os.path.join(self.output_dir, filename)

        # ── 1차 패스: BytesIO에 렌더링, 챕터 시작 페이지 수집 (이미지 생략)
        self.chapter_pages = {}
        buf = _io.BytesIO()
        c_pass1 = canvas.Canvas(buf, pagesize=A4)
        c_pass1.setTitle(_safe_text(title))
        self.y = self.page_height - self.margin_top
        self._page_drawn = False
        self._render_pages(c_pass1, ebook_data, skip_images=True)
        c_pass1.save()

        # ── 2차 패스: 실제 파일에 렌더링, 1차 패스에서 수집한 챕터 페이지 번호 사용
        c = canvas.Canvas(filepath, pagesize=A4)
        c.setTitle(_safe_text(title))
        self.y = self.page_height - self.margin_top
        self._page_drawn = False
        self._render_pages(c, ebook_data, skip_images=False)
        c.save()

        return filepath, filename

    # ──────────────────────────────────────────────
    # 표지 (아름다운 프로그래매틱 디자인)
    # ──────────────────────────────────────────────
    def _draw_cover(self, c, ebook_data):
        book_info = ebook_data.get('book_info', {})
        title = _safe_text(book_info.get('book_title', '전자책'))
        subtitle = _safe_text(book_info.get('subtitle', ''))
        W, H = self.page_width, self.page_height

        # ── 배경: 짙은 네이비 → 딥 퍼플 그라데이션 (레이어드 직사각형) ──
        bg_bands = [
            (0.00, 0.40, '#0a0a1a'),
            (0.40, 0.65, '#0e0e22'),
            (0.65, 0.85, '#12122e'),
            (0.85, 1.00, '#16163a'),
        ]
        for y_s, y_e, col in bg_bands:
            bh = H * (y_e - y_s)
            by = H * (1.0 - y_e)
            c.setFillColor(HexColor(col))
            c.rect(0, by, W, bh, fill=1, stroke=0)

        # ── 장식: 우상단 대형 원(겹침) ──
        c.setFillColor(HexColor('#1a1a40'))
        c.circle(W + 20, H - 30, 180, fill=1, stroke=0)
        c.setFillColor(HexColor('#20205a'))
        c.circle(W - 20, H - 60, 120, fill=1, stroke=0)
        c.setStrokeColor(HexColor('#6c5ce7'))
        c.setLineWidth(1.2)
        c.circle(W - 10, H - 50, 140, fill=0, stroke=1)

        # ── 장식: 좌하단 소형 원 ──
        c.setFillColor(HexColor('#141430'))
        c.circle(-30, 80, 100, fill=1, stroke=0)
        c.setStrokeColor(HexColor('#4a4a8a'))
        c.setLineWidth(0.8)
        c.circle(-20, 70, 120, fill=0, stroke=1)

        # ── 좌측 세로 액센트 바 ──
        c.setFillColor(HexColor('#6c5ce7'))
        c.rect(0, 0, 6, H, fill=1, stroke=0)

        # ── 상단 액센트 바 ──
        c.setFillColor(HexColor('#6c5ce7'))
        c.rect(0, H - 6, W, 6, fill=1, stroke=0)

        # ── 수평 장식선 (중간) ──
        line_y = H * 0.52
        c.setStrokeColor(HexColor('#6c5ce7'))
        c.setLineWidth(1.5)
        c.line(40, line_y, 40 + 60, line_y)
        c.setStrokeColor(HexColor('#4a4a6a'))
        c.setLineWidth(0.5)
        c.line(40 + 70, line_y, W - 40, line_y)

        # ── 카테고리 레이블 ──
        c.setFont(self.font_name, 8)
        c.setFillColor(HexColor('#9b8fff'))
        c.drawString(42, line_y + 14, '전자책  /  E-BOOK')

        # ── 제목 ──
        title_font_size = 28 if len(title) <= 14 else (22 if len(title) <= 22 else 18)
        title_y = line_y - 20
        title_lines = self._wrap_text_by_width(c, title, title_font_size, W - 90)
        for tl in title_lines:
            c.setFont(self.font_name, title_font_size)
            c.setFillColor(HexColor('#FFFFFF'))
            c.drawString(42, title_y, tl)
            title_y -= title_font_size * 1.45

        # ── 부제목 ──
        if subtitle:
            title_y -= 12
            sub_lines = self._wrap_text_by_width(c, subtitle, 13, W - 90)
            for sl in sub_lines:
                c.setFont(self.font_name, 13)
                c.setFillColor(HexColor('#aaaacc'))
                c.drawString(42, title_y, sl)
                title_y -= 13 * 1.5

        # ── 하단 장식선 ──
        c.setStrokeColor(HexColor('#6c5ce7'))
        c.setLineWidth(1.5)
        c.line(42, 55, 42 + 80, 55)

        # ── 하단 서브텍스트 ──
        c.setFont(self.font_name, 8)
        c.setFillColor(HexColor('#555577'))
        c.drawString(42, 38, 'AI 전자책 자동 생성기 제작')

        # ── 하단 가로 바 ──
        c.setFillColor(HexColor('#6c5ce7'))
        c.rect(0, 0, W, 5, fill=1, stroke=0)

    # ──────────────────────────────────────────────
    # 목차 (고급 디자인)
    # ──────────────────────────────────────────────
    def _draw_toc(self, c, ebook_data):
        # 헤더 배경 바
        c.setFillColor(HexColor('#1a1a2e'))
        c.rect(self.margin_left - 10, self.y - 8, self.content_width + 20, 38, fill=1, stroke=0)
        c.setFont(self.font_name, 18)
        c.setFillColor(HexColor('#FFFFFF'))
        c.drawString(self.margin_left + 6, self.y + 6, '목  차')
        self.y -= 46
        self._mark_drawn()

        # 구분선
        c.setStrokeColor(HexColor('#6c5ce7'))
        c.setLineWidth(2)
        c.line(self.margin_left, self.y, self.margin_left + self.content_width, self.y)
        self.y -= 16

        chapters = ebook_data.get('book_info', {}).get('chapters', [])
        current_phase = None
        has_page_nums = bool(self.chapter_pages)

        for i, ch in enumerate(chapters):
            self._check_page_break(c, 40)
            phase = ch.get('phase', '')
            num = ch.get('chapter_num', i + 1)
            ch_title = _safe_text(ch.get('title', ''))
            color = PHASE_COLORS.get(phase, '#6c5ce7')

            # 단계 변경 시 그룹 헤더
            if phase != current_phase:
                current_phase = phase
                if i > 0:
                    self.y -= 8
                c.setFont(self.font_name, 8)
                c.setFillColor(HexColor(color))
                c.drawString(self.margin_left, self.y, f'▶  {phase}')
                self.y -= 16
                c.setStrokeColor(HexColor('#eeeeee'))
                c.setLineWidth(0.5)
                c.line(self.margin_left, self.y + 2, self.margin_left + self.content_width, self.y + 2)
                self.y -= 6
                self._mark_drawn()

            self._check_page_break(c, 28)

            # 챕터 번호 원형 배지
            badge_x = self.margin_left + 10
            c.setFillColor(HexColor(color))
            c.circle(badge_x, self.y + 5, 8, fill=1, stroke=0)
            c.setFont(self.font_name, 7)
            c.setFillColor(HexColor('#FFFFFF'))
            c.drawCentredString(badge_x, self.y + 2, str(num))

            # 페이지 번호 (오른쪽)
            pg_text = str(self.chapter_pages.get(num, '')) if has_page_nums else ''
            pg_width = c.stringWidth(pg_text, self.font_name, 10) + 4 if pg_text else 0

            # 챕터 제목 (제목 영역 = 전체 - 배지 - 페이지번호)
            title_x = self.margin_left + 24
            available = self.content_width - 26 - pg_width
            title_lines = self._wrap_text_by_width(c, ch_title, 10, available)
            display = title_lines[0] if title_lines else ch_title
            if len(title_lines) > 1:
                display += '…'
            c.setFont(self.font_name, 10)
            c.setFillColor(HexColor('#333333'))
            c.drawString(title_x, self.y, display)

            # 페이지 번호 우측 정렬
            if pg_text:
                pg_x = self.margin_left + self.content_width - pg_width
                # 점선 리더
                title_end_x = title_x + c.stringWidth(display, self.font_name, 10) + 4
                dot_y = self.y + 3
                c.setFont(self.font_name, 8)
                c.setFillColor(HexColor('#cccccc'))
                dot_x = title_end_x
                while dot_x + 6 < pg_x - 4:
                    c.drawString(dot_x, dot_y, '.')
                    dot_x += 5
                c.setFont(self.font_name, 10)
                c.setFillColor(HexColor('#555555'))
                c.drawString(pg_x, self.y, pg_text)

            self.y -= 22
            self._mark_drawn()

        self.y -= 10

    # ──────────────────────────────────────────────
    # 가치 요약
    # ──────────────────────────────────────────────
    def _draw_analysis_summary(self, c, ebook_data):
        analysis = ebook_data.get('analysis', {})
        if not analysis:
            return

        # 섹션 헤더
        c.setFillColor(HexColor('#f8f7ff'))
        c.rect(self.margin_left - 10, self.y - 8, self.content_width + 20, 38, fill=1, stroke=0)
        c.setStrokeColor(HexColor('#6c5ce7'))
        c.setLineWidth(3)
        c.line(self.margin_left - 10, self.y - 8, self.margin_left - 10, self.y + 30)
        c.setFont(self.font_name, 18)
        c.setFillColor(HexColor('#1a1a2e'))
        c.drawString(self.margin_left + 6, self.y + 6, '이 책이 주는 가치')
        self.y -= 52
        self._mark_drawn()

        problem = analysis.get('problem_solved', {})
        items = [
            ('시간 절약',  problem.get('time', ''),    '#03c75a'),
            ('비용 절감',  problem.get('money', ''),   '#f39c12'),
            ('감정적 해방', problem.get('emotion', ''), '#3498db'),
        ]
        for label, val, col in items:
            if not val:
                continue
            self._check_page_break(c, 60)
            # 라벨 배지
            badge_w = c.stringWidth(label, self.font_name, 9) + 14
            c.setFillColor(HexColor(col))
            c.roundRect(self.margin_left, self.y - 2, badge_w, 16, 4, fill=1, stroke=0)
            c.setFont(self.font_name, 9)
            c.setFillColor(HexColor('#FFFFFF'))
            c.drawString(self.margin_left + 7, self.y + 1, label)
            self.y -= 20
            self._draw_paragraph(c, val, indent=10)
            self.y -= 8
            self._mark_drawn()

        why_pay = analysis.get('why_pay', '')
        if why_pay:
            self._check_page_break(c, 60)
            self.y -= 10
            c.setStrokeColor(HexColor('#6c5ce7'))
            c.setLineWidth(2)
            c.line(self.margin_left, self.y, self.margin_left + 40, self.y)
            self.y -= 16
            c.setFont(self.font_name, 13)
            c.setFillColor(HexColor('#1a1a2e'))
            c.drawString(self.margin_left, self.y, '왜 이 책에 투자해야 하는가')
            self.y -= 20
            self._draw_paragraph(c, why_pay, indent=10)
            self._mark_drawn()

    # ──────────────────────────────────────────────
    # 프롤로그
    # ──────────────────────────────────────────────
    def _draw_prologue(self, c, prologue_text):
        # 헤더
        c.setFillColor(HexColor('#1a1a2e'))
        c.rect(self.margin_left - 10, self.y - 8, self.content_width + 20, 38, fill=1, stroke=0)
        c.setFont(self.font_name, 18)
        c.setFillColor(HexColor('#FFFFFF'))
        c.drawString(self.margin_left + 6, self.y + 6, '프롤로그')
        self.y -= 52
        c.setStrokeColor(HexColor('#6c5ce7'))
        c.setLineWidth(2)
        c.line(self.margin_left, self.y, self.margin_left + self.content_width, self.y)
        self.y -= 20
        self._mark_drawn()
        self._draw_paragraph(c, prologue_text, font_size=self.font_size, color='#333333')

    # ──────────────────────────────────────────────
    # 에필로그
    # ──────────────────────────────────────────────
    def _draw_epilogue(self, c, epilogue_text):
        # 헤더
        c.setFillColor(HexColor('#0a0a1a'))
        c.rect(self.margin_left - 10, self.y - 8, self.content_width + 20, 38, fill=1, stroke=0)
        c.setFont(self.font_name, 18)
        c.setFillColor(HexColor('#FFFFFF'))
        c.drawString(self.margin_left + 6, self.y + 6, '에필로그')
        self.y -= 52
        c.setStrokeColor(HexColor('#6c5ce7'))
        c.setLineWidth(2)
        c.line(self.margin_left, self.y, self.margin_left + self.content_width, self.y)
        self.y -= 20
        self._mark_drawn()
        self._draw_paragraph(c, epilogue_text, font_size=self.font_size, color='#333333')

    # ──────────────────────────────────────────────
    # 챕터 시작 페이지 (고급 디자인)
    # ──────────────────────────────────────────────
    def _draw_chapter_start(self, c, chapter, num, img_url=None):
        phase = chapter.get('phase', '')
        title = _safe_text(chapter.get('title', ''))
        before = _safe_text(chapter.get('before_state', ''))
        after = _safe_text(chapter.get('after_state', ''))
        color = PHASE_COLORS.get(phase, '#6c5ce7')
        bg_color = PHASE_BG.get(phase, '#f8f7ff')
        W = self.page_width

        # 상단 컬러 헤더 바
        c.setFillColor(HexColor(color))
        c.rect(0, self.page_height - self.margin_top - 45 + self.heading_size * 1.4 + 30,
               W, 6, fill=1, stroke=0)

        # 단계 레이블
        phase_w = c.stringWidth(f'  {phase}  ', self.font_name, 9) + 4
        c.setFillColor(HexColor(color))
        c.roundRect(self.margin_left, self.y - 2, phase_w, 18, 3, fill=1, stroke=0)
        c.setFont(self.font_name, 9)
        c.setFillColor(HexColor('#FFFFFF'))
        c.drawString(self.margin_left + 7, self.y + 2, phase)
        self.y -= 26
        self._mark_drawn()

        # CHAPTER N
        c.setFont(self.font_name, 11)
        c.setFillColor(HexColor('#aaaaaa'))
        c.drawString(self.margin_left, self.y, f'CHAPTER  {num}')
        self.y -= 22

        # 챕터 제목 (대형)
        title_lines = self._wrap_text_by_width(c, title, self.heading_size, self.content_width)
        c.setFillColor(HexColor('#1a1a1a'))
        for tl in title_lines:
            self._check_page_break(c, self.heading_size * 1.5)
            c.setFont(self.font_name, self.heading_size)
            c.setFillColor(HexColor('#1a1a1a'))
            c.drawString(self.margin_left, self.y, tl)
            self.y -= self.heading_size * 1.4
            self._mark_drawn()
        self.y -= 8

        # 밑줄
        c.setStrokeColor(HexColor(color))
        c.setLineWidth(2.5)
        c.line(self.margin_left, self.y, self.margin_left + 55, self.y)
        c.setStrokeColor(HexColor('#eeeeee'))
        c.setLineWidth(0.8)
        c.line(self.margin_left + 65, self.y, self.margin_left + self.content_width, self.y)
        self.y -= 20

        # 챕터 이미지 (다운로드 성공 시)
        if img_url:
            img_path = download_image(img_url, self.output_dir)
            if img_path:
                try:
                    img_height = min(160, self.y - self.margin_bottom - 80)
                    if img_height > 60:
                        c.drawImage(img_path, self.margin_left, self.y - img_height,
                                    width=self.content_width, height=img_height,
                                    preserveAspectRatio=True, anchor='c')
                        self.y -= img_height + 15
                        self._mark_drawn()
                except Exception:
                    pass

        # Before / After 박스
        if before or after:
            box_h = 14 + 14 * 2.2 * 2  # 예상 높이
            self._check_page_break(c, box_h + 20)

            # 박스 배경
            c.setFillColor(HexColor(bg_color))
            c.roundRect(self.margin_left, self.y - box_h + 10,
                        self.content_width, box_h, 6, fill=1, stroke=0)
            c.setStrokeColor(HexColor(color))
            c.setLineWidth(1)
            c.roundRect(self.margin_left, self.y - box_h + 10,
                        self.content_width, box_h, 6, fill=0, stroke=1)
            self.y -= 10

            if before:
                c.setFont(self.font_name, 8)
                c.setFillColor(HexColor('#e74c3c'))
                c.drawString(self.margin_left + 10, self.y, '읽기 전')
                self.y -= 14
                before_lines = self._wrap_text_by_width(c, before, 10, self.content_width - 20)
                for bl in before_lines[:2]:
                    c.setFont(self.font_name, 10)
                    c.setFillColor(HexColor('#cc3333'))
                    c.drawString(self.margin_left + 14, self.y, bl)
                    self.y -= 14
                    self._mark_drawn()
                self.y -= 4

            if after:
                c.setFont(self.font_name, 8)
                c.setFillColor(HexColor('#03c75a'))
                c.drawString(self.margin_left + 10, self.y, '읽고 난 후')
                self.y -= 14
                after_lines = self._wrap_text_by_width(c, after, 10, self.content_width - 20)
                for al in after_lines[:2]:
                    c.setFont(self.font_name, 10)
                    c.setFillColor(HexColor('#017a38'))
                    c.drawString(self.margin_left + 14, self.y, al)
                    self.y -= 14
                    self._mark_drawn()

    # ──────────────────────────────────────────────
    # 챕터 본문
    # ──────────────────────────────────────────────
    def _draw_chapter_content(self, c, content, chapter_num):
        content = _safe_text(content)
        if not content.strip():
            return

        lines = content.split('\n')
        for line in lines:
            stripped = line.strip()
            if not stripped:
                self.y -= self.font_size * self.line_spacing * 0.4
                continue

            # === 소제목 === 패턴
            submatch = re.match(r'^={2,}\s*(.+?)\s*={2,}$', stripped)
            if submatch:
                self._check_page_break(c, 55)
                self.y -= 12
                sub_title = _safe_text(submatch.group(1))

                # 소제목 배경 바
                sub_bg_h = self.subheading_size * 1.6
                c.setFillColor(HexColor('#f4f4f8'))
                c.rect(self.margin_left - 6, self.y - sub_bg_h + self.subheading_size,
                       self.content_width + 12, sub_bg_h, fill=1, stroke=0)
                # 좌측 강조선
                c.setFillColor(HexColor('#6c5ce7'))
                c.rect(self.margin_left - 6, self.y - sub_bg_h + self.subheading_size,
                       3, sub_bg_h, fill=1, stroke=0)

                c.setFont(self.font_name, self.subheading_size)
                c.setFillColor(HexColor('#1a1a2e'))
                c.drawString(self.margin_left + 4, self.y, sub_title)
                self.y -= self.subheading_size * 1.7
                self._mark_drawn()
                continue

            # [핵심 포인트] / [실전 팁] 박스
            tip_match = re.match(r'^\[(핵심 포인트|실전 팁|TIP|핵심|포인트|POINT)\]\s*(.*)', stripped)
            if tip_match:
                tip_text = _safe_text(stripped)
                self._check_page_break(c, 40)
                self.y -= 4
                # 배경 박스
                est_lines = max(1, len(tip_text) // 40)
                box_h = (self.font_size * self.line_spacing) * est_lines + 16
                c.setFillColor(HexColor('#f0fff4'))
                c.roundRect(self.margin_left, self.y - box_h + self.font_size + 4,
                            self.content_width, box_h, 5, fill=1, stroke=0)
                c.setStrokeColor(HexColor('#03c75a'))
                c.setLineWidth(1)
                c.roundRect(self.margin_left, self.y - box_h + self.font_size + 4,
                            self.content_width, box_h, 5, fill=0, stroke=1)
                self._draw_paragraph(c, tip_text, font_size=self.font_size,
                                     color='#1a6b35', indent=8)
                self.y -= 6
                continue

            # 일반 문단
            self._draw_paragraph(c, stripped)

    # ──────────────────────────────────────────────
    # 마케팅 부록
    # ──────────────────────────────────────────────
    def _draw_marketing_page(self, c, ebook_data):
        marketing = ebook_data.get('marketing', {})
        if not marketing:
            return

        # 헤더
        c.setFillColor(HexColor('#1a1a2e'))
        c.rect(self.margin_left - 10, self.y - 8, self.content_width + 20, 38, fill=1, stroke=0)
        c.setFont(self.font_name, 16)
        c.setFillColor(HexColor('#FFFFFF'))
        c.drawString(self.margin_left + 6, self.y + 6, '부록 : 이 책에 대하여')
        self.y -= 50
        self._mark_drawn()

        sales_copy = marketing.get('sales_copy', '')
        if sales_copy:
            c.setFont(self.font_name, 12)
            c.setFillColor(HexColor('#03c75a'))
            c.drawString(self.margin_left, self.y, '판매 소개문')
            self.y -= 18

            # 인용문 박스
            c.setFillColor(HexColor('#f8f7ff'))
            c.roundRect(self.margin_left, self.y - 10,
                        self.content_width, 12, 4, fill=1, stroke=0)
            c.setStrokeColor(HexColor('#6c5ce7'))
            c.setLineWidth(2)
            c.line(self.margin_left, self.y + 12,
                   self.margin_left, self.y - self.font_size * self.line_spacing * 8)
            self._draw_paragraph(c, sales_copy, indent=12)
            self.y -= 20

        value = marketing.get('value_summary', {})
        if value:
            c.setFont(self.font_name, 12)
            c.setFillColor(HexColor('#03c75a'))
            c.drawString(self.margin_left, self.y, '독자에게 주는 가치')
            self.y -= 20
            self._mark_drawn()
            for key, label in [('time_saved', '시간 절약'), ('money_saved', '비용 절감'), ('mistakes_prevented', '실수 방지')]:
                val = value.get(key, '')
                if val:
                    self._check_page_break(c, 30)
                    self._draw_paragraph(c, f'• {label}: {val}', font_size=10)


# ============================================================
# DOCX 생성기
# ============================================================
class EbookDocxGenerator:
    def __init__(self, config=None):
        self.config = config or load_config()
        self.output_dir = self.config.get('output_dir', './static/output')
        os.makedirs(self.output_dir, exist_ok=True)

    def generate(self, ebook_data):
        from docx import Document
        from docx.shared import Pt, Cm, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        book_info = ebook_data.get('book_info', {})
        title = book_info.get('book_title', '전자책')
        safe_title = re.sub(r'[^\w가-힣\s-]', '', title)[:50].strip()
        filename = f"{safe_title}.docx"
        filepath = os.path.join(self.output_dir, filename)

        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        doc = Document()
        section = doc.sections[0]
        section.page_width = Cm(21.0)
        section.page_height = Cm(29.7)
        ml = self.config.get('pdf_margin_left', 60) / 72 * 2.54
        mr = self.config.get('pdf_margin_right', 60) / 72 * 2.54
        mt = self.config.get('pdf_margin_top', 72) / 72 * 2.54
        mb = self.config.get('pdf_margin_bottom', 72) / 72 * 2.54
        section.left_margin   = Cm(ml)
        section.right_margin  = Cm(mr)
        section.top_margin    = Cm(mt)
        section.bottom_margin = Cm(mb)

        # 페이지 번호 푸터 추가
        footer = section.footer
        fp = footer.paragraphs[0]
        fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        def _add_page_num_field(para):
            run = para.add_run('— ')
            run.font.size = Pt(9); run.font.color.rgb = RGBColor(0xaa,0xaa,0xaa)
            for tag, text in [('w:fldChar','begin'),('w:instrText','PAGE'),('w:fldChar','end')]:
                el = OxmlElement(tag)
                if tag == 'w:instrText':
                    el.text = text
                    run2 = para.add_run(); run2._r.append(el)
                else:
                    el.set(qn('w:fldCharType'), text)
                    run2 = para.add_run(); run2._r.append(el)
            run3 = para.add_run(' —')
            run3.font.size = Pt(9); run3.font.color.rgb = RGBColor(0xaa,0xaa,0xaa)
        _add_page_num_field(fp)

        fs = self.config.get('pdf_font_size', 11)
        hs = self.config.get('pdf_heading_size', 16)
        ss = self.config.get('pdf_subheading_size', 13)
        ls = self.config.get('pdf_line_spacing', 1.6)

        # 표지
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.add_run('\n\n\n\n')
        r = p.add_run(title)
        r.bold = True; r.font.size = Pt(28)
        r.font.color.rgb = RGBColor(0x1a, 0x1a, 0x2e)
        subtitle = book_info.get('subtitle', '')
        if subtitle:
            p2 = doc.add_paragraph()
            p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
            r2 = p2.add_run(subtitle)
            r2.font.size = Pt(14)
            r2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        doc.add_page_break()

        # 목차 (추정 페이지번호 포함)
        h = doc.add_heading('목차', level=1)
        for run in h.runs: run.font.size = Pt(22)

        # ── 페이지번호 추정 ──────────────────────────────────────
        # A4 기준: 본문 영역 높이에서 행 수 × 폰트 크기로 1페이지 글자수 추정
        # 한국어 A4 기준 약 600~700자/페이지 (여백, 줄간격 고려)
        CHARS_PER_PAGE = 600
        est_page = 2  # 표지(1p) + 목차 시작(2p)
        # 목차 자체가 차지하는 페이지 (챕터 수에 따라)
        n_chapters = len(book_info.get('chapters', []))
        est_page += max(1, (n_chapters + 15) // 16)  # ~16 항목/페이지

        # 프롤로그
        prologue_text = ebook_data.get('prologue', '')
        if prologue_text and prologue_text.strip():
            est_page += max(1, len(prologue_text) // CHARS_PER_PAGE)

        # 분석/가치 요약
        analysis_text = ebook_data.get('analysis', {})
        if analysis_text:
            est_page += 1

        # 각 챕터 시작 페이지 추정
        chapter_est_pages = {}
        for i, ch_data in enumerate(ebook_data.get('chapters_content', [])):
            ch_num = book_info.get('chapters', [{}])[i].get('chapter_num', i + 1) if i < n_chapters else i + 1
            chapter_est_pages[ch_num] = est_page
            content = ch_data.get('content', '')
            # 챕터 내용 길이 기반 페이지 추정
            est_page += max(1, len(content) // CHARS_PER_PAGE)

        # 목차 항목에 탭 + 추정 페이지 번호 추가
        for ch in book_info.get('chapters', []):
            ch_num = ch.get('chapter_num', '')
            p = doc.add_paragraph()
            # 탭 스톱 설정 (우측 정렬, 점선 리더)
            pPr = p._p.get_or_add_pPr()
            tabs_el = OxmlElement('w:tabs')
            tab_el = OxmlElement('w:tab')
            tab_el.set(qn('w:val'), 'right')
            tab_el.set(qn('w:leader'), 'dot')
            # A4 폭(21cm) - 좌우여백으로 탭 위치 계산 (twips: 1cm = 567 twips)
            tab_pos = int((21.0 - ml - mr) * 567)
            tab_el.set(qn('w:pos'), str(tab_pos))
            tabs_el.append(tab_el)
            pPr.append(tabs_el)

            r1 = p.add_run(f"[{ch.get('phase','')}]  ")
            r1.font.size = Pt(9); r1.font.color.rgb = RGBColor(0x88,0x88,0x88)
            r2 = p.add_run(f"CH.{ch_num}  {ch.get('title','')}")
            r2.font.size = Pt(fs)
            # 탭 + 추정 페이지 번호 (plain text)
            r_tab = p.add_run('\t')
            r_tab.font.size = Pt(fs)
            pg_num = chapter_est_pages.get(ch_num, '')
            r_pg = p.add_run(str(pg_num))
            r_pg.font.size = Pt(fs)
            r_pg.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        doc.add_page_break()

        # 프롤로그
        prologue = ebook_data.get('prologue', '')
        if prologue and prologue.strip():
            h = doc.add_heading('프롤로그', level=1)
            for run in h.runs: run.font.size = Pt(20)
            p = doc.add_paragraph(prologue)
            p.paragraph_format.line_spacing = ls
            for run in p.runs: run.font.size = Pt(fs)
            doc.add_page_break()

        # 가치 요약
        analysis = ebook_data.get('analysis', {})
        if analysis:
            h = doc.add_heading('이 책이 주는 가치', level=1)
            problem = analysis.get('problem_solved', {})
            for label, key in [('시간 절약','time'),('비용 절감','money'),('감정적 해방','emotion')]:
                doc.add_heading(label, level=3)
                p = doc.add_paragraph(problem.get(key,''))
                p.paragraph_format.line_spacing = ls
                for run in p.runs: run.font.size = Pt(fs)
            if analysis.get('why_pay'):
                doc.add_heading('왜 이 책에 투자해야 하는가', level=2)
                p = doc.add_paragraph(analysis['why_pay'])
                p.paragraph_format.line_spacing = ls
                for run in p.runs: run.font.size = Pt(fs)
            doc.add_page_break()

        # 챕터 (각 챕터에 북마크 삽입 — 목차 PAGEREF 연동)
        def _add_bookmark(paragraph, bm_name, bm_id):
            """단락에 북마크(bookmarkStart + bookmarkEnd)를 추가"""
            bm_start = OxmlElement('w:bookmarkStart')
            bm_start.set(qn('w:id'), str(bm_id))
            bm_start.set(qn('w:name'), bm_name)
            paragraph._p.insert(0, bm_start)
            bm_end = OxmlElement('w:bookmarkEnd')
            bm_end.set(qn('w:id'), str(bm_id))
            paragraph._p.append(bm_end)

        for i, ch_data in enumerate(ebook_data.get('chapters_content', [])):
            chapter = ch_data.get('chapter', {})
            content = ch_data.get('content', '')

            h = doc.add_heading(f"CHAPTER {i+1}", level=2)
            for run in h.runs: run.font.size = Pt(12); run.font.color.rgb = RGBColor(0xaa,0xaa,0xaa)
            # 챕터 제목에 북마크 추가 (목차 페이지 번호 연동)
            h2 = doc.add_heading(chapter.get('title',''), level=1)
            for run in h2.runs: run.font.size = Pt(hs)
            _add_bookmark(h2, f'_ch{i+1}', i + 100)

            before = chapter.get('before_state','')
            after  = chapter.get('after_state','')
            if before:
                p = doc.add_paragraph()
                r = p.add_run(f"읽기 전: {before}")
                r.font.size = Pt(10); r.font.color.rgb = RGBColor(0xe7,0x4c,0x3c)
            if after:
                p = doc.add_paragraph()
                r = p.add_run(f"읽고 난 후: {after}")
                r.font.size = Pt(10); r.font.color.rgb = RGBColor(0x03,0xc7,0x5a)
            doc.add_paragraph('')

            for line in content.split('\n'):
                stripped = line.strip()
                if not stripped:
                    doc.add_paragraph('')
                    continue
                submatch = re.match(r'^={2,}\s*(.+?)\s*={2,}$', stripped)
                if submatch:
                    h = doc.add_heading(submatch.group(1), level=2)
                    for run in h.runs: run.font.size = Pt(ss)
                    continue
                p = doc.add_paragraph(stripped)
                p.paragraph_format.line_spacing = ls
                for run in p.runs: run.font.size = Pt(fs)

            doc.add_page_break()

        # 에필로그
        epilogue = ebook_data.get('epilogue', '')
        if epilogue and epilogue.strip():
            h = doc.add_heading('에필로그', level=1)
            for run in h.runs: run.font.size = Pt(20)
            p = doc.add_paragraph(epilogue)
            p.paragraph_format.line_spacing = ls
            for run in p.runs: run.font.size = Pt(fs)
            doc.add_page_break()

        # 마케팅
        marketing = ebook_data.get('marketing', {})
        if marketing:
            doc.add_heading('부록: 이 책에 대하여', level=1)
            if marketing.get('sales_copy'):
                doc.add_heading('판매 소개문', level=2)
                p = doc.add_paragraph(marketing['sales_copy'])
                p.paragraph_format.line_spacing = ls
                for run in p.runs: run.font.size = Pt(fs)

        doc.save(filepath)
        return filepath, filename


# ============================================================
# PPTX 생성기  (v3 — 완결 문장, 조화로운 레이아웃)
# ============================================================
class EbookPptxGenerator:
    """파워포인트 스타일 PPTX 생성기 v3 — 완결 문장, 조화로운 레이아웃"""

    SLIDE_W  = 13.33   # 16:9 와이드스크린 인치
    SLIDE_H  = 7.5
    HEADER_H = 1.15    # 표준 헤더 높이
    MX       = 0.55    # 좌우 기본 여백
    CONTENT_TOP = 1.3  # 컨텐츠 시작 y

    def __init__(self, config=None):
        self.config = config or load_config()
        self.output_dir = self.config.get('output_dir', './static/output')
        os.makedirs(self.output_dir, exist_ok=True)

    # ── primitives ────────────────────────────────────────────────────
    @staticmethod
    def _rgb(hex_color):
        from pptx.dml.color import RGBColor
        h = hex_color.lstrip('#')[:6]
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

    @staticmethod
    def _blank(prs):
        return prs.slides.add_slide(prs.slide_layouts[6])

    @staticmethod
    def _set_bg(slide, hex_color):
        from pptx.dml.color import RGBColor
        fill = slide.background.fill
        fill.solid()
        h = hex_color.lstrip('#')[:6]
        fill.fore_color.rgb = RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

    def _rect(self, slide, x, y, w, h, fill_hex):
        from pptx.util import Inches
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._rgb(fill_hex)
        shape.line.fill.background()
        return shape

    def _txt(self, slide, text, x, y, w, h,
             size=14, bold=False, color='#333333',
             align=None, wrap=True, italic=False):
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.text = str(text)
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = self._rgb(color)
        p.alignment = align or PP_ALIGN.LEFT
        return box

    # ── text utilities ────────────────────────────────────────────────
    def _natural_trim(self, text, max_chars):
        """한국어 자연 종결 경계에서 자르기 (문장 중간 절단 방지)"""
        if len(text) <= max_chars:
            return text
        # 종결 어미 탐색 (max_chars + 10 범위 내)
        endings = ['합니다', '됩니다', '습니다', '입니다', '하세요',
                   '이다', '한다', '된다', '있다', '없다', '이며', '하며', '하고']
        chunk = text[:max_chars + 10]
        best = 0
        for e in endings:
            p = chunk.rfind(e)
            if p > max_chars // 3:
                ep = p + len(e)
                if ep <= max_chars + 10 and ep > best:
                    best = ep
        if best > max_chars // 3:
            return text[:best]
        # 공백 경계
        p = text[:max_chars].rfind(' ')
        if p > max_chars // 2:
            return text[:p]
        return text[:max_chars]

    def _extract_bullets(self, lines, max_bullets=5, max_chars=95):
        """완결 문장/핵심 포인트 추출 — 문장 절단 없음"""
        bullets = []
        for line in lines:
            line = line.strip()
            if not line or len(line) < 4:
                continue
            # ① 기호 리스트 (-, •, ●, ▶, ►, ✓ 등)
            m = re.match(r'^[-•●▶►✓☑✔▷※]\s+(.+)', line)
            if m:
                b = self._natural_trim(m.group(1).strip(), max_chars)
                if b:
                    bullets.append(b)
                if len(bullets) >= max_bullets:
                    break
                continue
            # ② 숫자 리스트
            m2 = re.match(r'^\d+[.)]\s+(.+)', line)
            if m2:
                b = self._natural_trim(m2.group(1).strip(), max_chars)
                if b:
                    bullets.append(b)
                if len(bullets) >= max_bullets:
                    break
                continue
            # ③ 일반 문단 → 완결된 첫 문장 추출
            if len(line) >= 8:
                parts = re.split(r'(?<=[.!?。])\s+', line)
                for part in parts:
                    part = part.strip()
                    if len(part) >= 6:
                        bullets.append(self._natural_trim(part, max_chars))
                        break
            if len(bullets) >= max_bullets:
                break
        return bullets or ['핵심 내용을 확인하세요.']

    # ── content parser ────────────────────────────────────────────────
    def _parse_sections(self, content):
        """챕터 컨텐츠 → 섹션 리스트
        각 섹션: {'title': str|None, 'type': str, 'lines': [str]}
        type: intro / section / highlight / summary / checklist
        """
        sections = []
        cur = {'title': None, 'type': 'intro', 'lines': []}

        for raw in content.split('\n'):
            line = raw.strip()
            # === 제목 ===
            m = re.match(r'^={2,}\s*(.+?)\s*={2,}$', line)
            if m:
                if any(l for l in cur['lines'] if l.strip()):
                    sections.append(cur)
                title = m.group(1)
                if '핵심 요약' in title or title.strip() == '요약':
                    t = 'summary'
                elif '체크리스트' in title or ('실행' in title and '리스트' in title):
                    t = 'checklist'
                elif '핵심 포인트' in title or '실전 팁' in title:
                    t = 'highlight'
                else:
                    t = 'section'
                cur = {'title': title, 'type': t, 'lines': []}
                continue
            # [핵심 포인트] 인라인 박스
            bm = re.match(r'^\[([^\]]{2,20})\]\s*(.*)', line)
            if bm:
                label = bm.group(1)
                rest = bm.group(2).strip()
                if '핵심' in label or '팁' in label or '포인트' in label:
                    if any(l for l in cur['lines'] if l.strip()):
                        sections.append(cur)
                    cur = {'title': label, 'type': 'highlight',
                           'lines': [rest] if rest else []}
                    continue
            cur['lines'].append(line)

        if any(l for l in cur['lines'] if l.strip()):
            sections.append(cur)
        return sections

    # ── slide builders ─────────────────────────────────────────────────
    def _slide_cover(self, prs, title, subtitle, total_chapters):
        from pptx.enum.text import PP_ALIGN
        W, H = self.SLIDE_W, self.SLIDE_H
        slide = self._blank(prs)
        self._set_bg(slide, '#0d1b2a')

        # 왼쪽 강조 바
        self._rect(slide, 0, 0, 0.28, H, '#4361ee')
        # 상/하단 라인
        self._rect(slide, 0.28, 0.55, W - 0.28, 0.04, '#4361ee')
        self._rect(slide, 0.28, H - 0.6, W - 0.28, 0.04, '#4361ee')

        # 제목
        self._txt(slide, title, 0.65, 1.45, W - 1.05, 2.9,
                  size=38, bold=True, color='#FFFFFF', wrap=True)
        # 부제목
        if subtitle:
            self._txt(slide, subtitle, 0.65, 4.55, W - 1.05, 1.0,
                      size=18, color='#7a9cc0', wrap=True)
        # 챕터 수 배지
        self._rect(slide, 0.65, 5.85, 3.8, 0.58, '#1e3356')
        self._rect(slide, 0.65, 5.85, 0.08, 0.58, '#4361ee')
        self._txt(slide, f"  총 {total_chapters}개 챕터  ·  실전 완전 가이드",
                  0.67, 5.9, 3.75, 0.48, size=13, color='#90b8d8')

    def _slide_toc(self, prs, chapters, chapter_slide_nums=None):
        from pptx.enum.text import PP_ALIGN
        W, H = self.SLIDE_W, self.SLIDE_H
        chapter_slide_nums = chapter_slide_nums or {}
        # 챕터 수에 따라 1~2 슬라이드 분할
        half = (len(chapters) + 1) // 2
        for page_idx, chunk in enumerate([chapters[:half], chapters[half:]]):
            if not chunk:
                continue
            slide = self._blank(prs)
            self._set_bg(slide, '#f4f6fb')
            # 헤더
            self._rect(slide, 0, 0, W, 1.05, '#0d1b2a')
            self._txt(slide, '목  차' if page_idx == 0 else '목  차 (계속)',
                      self.MX, 0.18, 5.0, 0.7, size=26, bold=True, color='#FFFFFF')
            self._txt(slide, 'Contents', W - 2.6, 0.28, 2.3, 0.55,
                      size=13, color='#667799', align=PP_ALIGN.RIGHT)

            y = 1.15
            n = len(chunk)
            row_h = min((H - 1.3) / max(n, 1), 0.73)
            for j, ch in enumerate(chunk):
                phase = ch.get('phase', '')
                pcol = PHASE_COLORS.get(phase, '#6c5ce7')
                num = ch.get('chapter_num', j + 1 + page_idx * half)
                ch_title = ch.get('title', '')

                # 짝수 행 배경
                bg = '#eaeef5' if j % 2 == 0 else '#f4f6fb'
                self._rect(slide, self.MX - 0.05, y - 0.03, W - self.MX * 2 + 0.1, row_h, bg)
                # 번호 배지
                self._rect(slide, self.MX - 0.05, y - 0.03, 0.55, row_h, pcol)
                self._txt(slide, str(num), self.MX - 0.05, y + 0.06,
                          0.55, row_h - 0.12,
                          size=13, bold=True, color='#FFFFFF', align=PP_ALIGN.CENTER)
                # 챕터 제목
                self._txt(slide, ch_title[:60], self.MX + 0.6, y,
                          W - self.MX * 2 - 2.4, row_h,
                          size=15, color='#1a1a2e', wrap=True)
                # 단계 태그
                if phase:
                    self._rect(slide, W - self.MX - 1.8, y + (row_h - 0.38) / 2,
                               1.05, 0.38, '#dde3ef')
                    self._txt(slide, phase, W - self.MX - 1.8, y + (row_h - 0.38) / 2,
                              1.05, 0.38, size=10, bold=True, color=pcol,
                              align=PP_ALIGN.CENTER)
                # 슬라이드 번호 (우측)
                pg = chapter_slide_nums.get(num, '')
                if pg:
                    self._txt(slide, str(pg), W - self.MX - 0.6, y + (row_h - 0.38) / 2,
                              0.55, 0.38, size=12, bold=True, color='#667799',
                              align=PP_ALIGN.CENTER)
                y += row_h + 0.03

    def _slide_chapter_intro(self, prs, chapter, ch_num, col_hex):
        from pptx.enum.text import PP_ALIGN
        W, H = self.SLIDE_W, self.SLIDE_H
        slide = self._blank(prs)

        bg = PHASE_BG.get(chapter.get('phase', ''), '#f4f6fb')
        self._set_bg(slide, bg)

        # 왼쪽 세로 바
        self._rect(slide, 0, 0, 0.22, H, col_hex)

        # 단계 배지
        phase_label = f"  {chapter.get('phase', '')}  ·  Chapter {ch_num:02d}"
        badge_w = max(3.0, len(phase_label) * 0.18)
        self._rect(slide, 0.45, 0.42, badge_w, 0.6, col_hex)
        self._txt(slide, phase_label, 0.47, 0.44, badge_w - 0.04, 0.56,
                  size=14, bold=True, color='#FFFFFF')

        # 챕터 제목
        self._txt(slide, chapter.get('title', ''), 0.45, 1.3, W - 0.85, 2.6,
                  size=32, bold=True, color='#1a1a2e', wrap=True)

        # 구분선
        self._rect(slide, 0.45, 4.1, W - 0.7, 0.04, col_hex)

        # Before / After
        before = chapter.get('before_state', '')
        after  = chapter.get('after_state', '')
        if before and after:
            # Before
            self._rect(slide, 0.45, 4.28, 6.0, 0.72, '#eeeeee')
            self._rect(slide, 0.45, 4.28, 0.08, 0.72, '#e74c3c')
            self._txt(slide, '읽기 전', 0.65, 4.3, 1.5, 0.32, size=10,
                      bold=True, color='#e74c3c')
            self._txt(slide, before[:60], 0.65, 4.56, 5.65, 0.4, size=12,
                      color='#555555', wrap=True)
            # After
            self._rect(slide, 0.45, 5.12, 6.0, 0.72, '#eaf7ef')
            self._rect(slide, 0.45, 5.12, 0.08, 0.72, col_hex)
            self._txt(slide, '읽고 난 후', 0.65, 5.14, 1.8, 0.32, size=10,
                      bold=True, color=col_hex)
            self._txt(slide, after[:60], 0.65, 5.4, 5.65, 0.4, size=12,
                      color='#1a1a2e', bold=True, wrap=True)
        elif before:
            self._rect(slide, 0.45, 4.28, 8.0, 0.68, '#eeeeee')
            self._txt(slide, f"Before → {before[:70]}", 0.65, 4.32, 7.8, 0.6,
                      size=12, color='#555555', wrap=True)
        elif after:
            self._rect(slide, 0.45, 4.28, 8.0, 0.68, '#eaf7ef')
            self._txt(slide, f"After → {after[:70]}", 0.65, 4.32, 7.8, 0.6,
                      size=12, color='#1a1a2e', bold=True, wrap=True)

        # 챕터 목적
        purpose = chapter.get('purpose', '')
        if purpose:
            self._txt(slide, f"✦  {purpose[:85]}", 0.45, 6.1, W - 0.7, 0.85,
                      size=13, color='#666666', italic=True, wrap=True)

    def _slide_section(self, prs, section_title, bullets, ch_label, col_hex):
        """일반 섹션 슬라이드 — 개별 textbox per bullet, 동적 폰트"""
        from pptx.enum.text import PP_ALIGN
        W, H = self.SLIDE_W, self.SLIDE_H
        slide = self._blank(prs)
        self._set_bg(slide, '#ffffff')

        # 헤더
        self._rect(slide, 0, 0, W, self.HEADER_H, col_hex)
        self._txt(slide, ch_label, self.MX, 0.1, W - self.MX - 0.5, 0.36,
                  size=10, color='#ffffffcc')
        self._txt(slide, section_title, self.MX, 0.42, W - self.MX - 0.4, 0.7,
                  size=22, bold=True, color='#FFFFFF', wrap=True)

        # 동적 폰트 크기
        n = len(bullets)
        total_chars = sum(len(b) for b in bullets)
        if n <= 2:
            fs = 22
        elif n == 3 and total_chars < 180:
            fs = 20
        elif n == 3:
            fs = 18
        elif n == 4 and total_chars < 240:
            fs = 18
        elif n == 4:
            fs = 16
        else:
            fs = 15

        # 컨텐츠 영역 — 균등 배분
        content_h = H - self.HEADER_H - 0.25
        item_h = content_h / max(n, 1)

        y = self.CONTENT_TOP
        for bullet in bullets:
            mid = y + item_h * 0.5 - 0.06  # 세로 중심
            # 왼쪽 강조 바 (세로 중심)
            accent_h = min(item_h * 0.6, 0.55)
            self._rect(slide, self.MX, mid - accent_h / 2, 0.07, accent_h, col_hex)
            # 불렛 텍스트 (개별 textbox — 높이 충분히)
            self._txt(slide, bullet, self.MX + 0.22, y,
                      W - self.MX - 0.55, item_h,
                      size=fs, color='#1a1a2e', wrap=True)
            y += item_h

    def _slide_highlight(self, prs, label, bullets, col_hex):
        """핵심 포인트 / 실전 팁 슬라이드"""
        from pptx.enum.text import PP_ALIGN
        W, H = self.SLIDE_W, self.SLIDE_H
        slide = self._blank(prs)
        self._set_bg(slide, '#f8f9ff')

        # 헤더 바
        self._rect(slide, 0, 0, W, self.HEADER_H, col_hex)
        self._txt(slide, f"★  {label}", self.MX, 0.26, W - self.MX * 2, 0.75,
                  size=24, bold=True, color='#FFFFFF')

        # 구분선
        self._rect(slide, self.MX, self.HEADER_H + 0.15, W - self.MX * 2, 0.04, col_hex)

        n = len(bullets)
        content_h = H - self.HEADER_H - 0.45
        item_h = content_h / max(n, 1)
        fs = 20 if n <= 2 else (18 if n <= 3 else 16)

        y = self.HEADER_H + 0.3
        for i, bullet in enumerate(bullets):
            # 번호 배지
            badge_size = 0.44
            badge_y = y + (item_h - badge_size) / 2
            self._rect(slide, self.MX, badge_y, badge_size, badge_size, col_hex)
            self._txt(slide, str(i + 1), self.MX, badge_y, badge_size, badge_size,
                      size=14, bold=True, color='#FFFFFF', align=PP_ALIGN.CENTER)
            # 내용
            self._txt(slide, bullet, self.MX + 0.58, y,
                      W - self.MX * 2 - 0.58, item_h,
                      size=fs, color='#1a1a2e', wrap=True)
            y += item_h

    def _slide_summary(self, prs, title, points, col_hex):
        """핵심 요약 — 카드 그리드"""
        from pptx.enum.text import PP_ALIGN
        W, H = self.SLIDE_W, self.SLIDE_H
        slide = self._blank(prs)
        self._set_bg(slide, '#0d1b2a')

        # 헤더
        self._rect(slide, 0, 0, W, self.HEADER_H, col_hex)
        self._txt(slide, f"핵심 정리  ·  {title}", self.MX, 0.24, W - self.MX * 2, 0.72,
                  size=22, bold=True, color='#FFFFFF')

        pts = points[:6]
        n = len(pts)
        cols = 3 if n >= 4 else min(n, 3)
        rows = (n + cols - 1) // cols
        pad = 0.12
        card_w = (W - self.MX * 2 - pad * (cols - 1)) / cols
        card_h = (H - self.HEADER_H - 0.3 - pad * (rows - 1)) / rows

        for idx, pt in enumerate(pts):
            row, col = divmod(idx, cols)
            cx = self.MX + col * (card_w + pad)
            cy = self.HEADER_H + 0.18 + row * (card_h + pad)
            self._rect(slide, cx, cy, card_w, card_h, '#1a2d40')
            self._rect(slide, cx, cy, 0.08, card_h, col_hex)
            # 번호
            self._txt(slide, str(idx + 1), cx + 0.14, cy + 0.06,
                      0.38, 0.38, size=13, bold=True, color=col_hex)
            # 내용 (자연 트림 적용)
            short = self._natural_trim(pt, 85)
            self._txt(slide, short, cx + 0.14, cy + 0.44,
                      card_w - 0.22, card_h - 0.52,
                      size=12, color='#b8d0e8', wrap=True)

    def _slide_checklist(self, prs, title, items, col_hex):
        """실행 체크리스트"""
        from pptx.enum.text import PP_ALIGN
        W, H = self.SLIDE_W, self.SLIDE_H
        slide = self._blank(prs)
        self._set_bg(slide, '#f0fff8')

        # 헤더
        self._rect(slide, 0, 0, W, self.HEADER_H, col_hex)
        self._txt(slide, f"실행 체크리스트  ·  {title}", self.MX, 0.24, W - self.MX * 2, 0.72,
                  size=21, bold=True, color='#FFFFFF')

        n = min(len(items), 7)
        content_h = H - self.HEADER_H - 0.25
        item_h = content_h / max(n, 1)
        fs = 18 if n <= 4 else (15 if n <= 6 else 13)

        y = self.CONTENT_TOP
        for idx, item in enumerate(items[:n]):
            bg = '#e3f5ec' if idx % 2 == 0 else '#f0fff8'
            self._rect(slide, self.MX - 0.05, y + 0.03, W - self.MX * 2 + 0.1, item_h - 0.06, bg)
            # 번호 배지
            badge_h = min(item_h - 0.12, 0.55)
            badge_y = y + (item_h - badge_h) / 2
            self._rect(slide, self.MX - 0.05, badge_y, 0.52, badge_h, col_hex)
            self._txt(slide, str(idx + 1), self.MX - 0.05, badge_y, 0.52, badge_h,
                      size=14, bold=True, color='#FFFFFF', align=PP_ALIGN.CENTER)
            # 내용
            self._txt(slide, item[:90], self.MX + 0.58, y,
                      W - self.MX * 2 - 0.58, item_h,
                      size=fs, color='#0a3a1a', wrap=True)
            y += item_h

    def _slide_text_page(self, prs, heading, body_text, bg_hex='#1a1a2e'):
        """프롤로그/에필로그용 텍스트 슬라이드"""
        from pptx.enum.text import PP_ALIGN
        W, H = self.SLIDE_W, self.SLIDE_H
        slide = self._blank(prs)
        self._set_bg(slide, bg_hex)
        # 제목 헤더
        self._rect(slide, 0, 0, W, self.HEADER_H, '#6c5ce7')
        self._txt(slide, heading, self.MX, 0.24, W - self.MX * 2, 0.72,
                  size=26, bold=True, color='#FFFFFF')
        # 본문 텍스트 (줄임)
        lines = [l.strip() for l in body_text.split('\n') if l.strip()]
        preview = ' '.join(lines)[:400] + ('…' if len(' '.join(lines)) > 400 else '')
        self._txt(slide, preview, self.MX, self.CONTENT_TOP,
                  W - self.MX * 2, H - self.CONTENT_TOP - 0.3,
                  size=15, color='#cccccc', wrap=True)

    def _add_slide_number(self, slide, num):
        """슬라이드 우측 하단에 페이지 번호 추가"""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        W, H = self.SLIDE_W, self.SLIDE_H
        txBox = slide.shapes.add_textbox(
            Inches(W - 0.7), Inches(H - 0.35), Inches(0.55), Inches(0.28)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(num)
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0xaa, 0xaa, 0xaa)

    # ── main ────────────────────────────────────────────────────────────
    def generate(self, ebook_data):
        from pptx import Presentation
        from pptx.util import Inches

        book_info = ebook_data.get('book_info', {})
        title    = book_info.get('book_title', '전자책')
        subtitle = book_info.get('subtitle', '')
        safe_title = re.sub(r'[^\w가-힣\s-]', '', title)[:50].strip()
        filename = f"{safe_title}.pptx"
        filepath = os.path.join(self.output_dir, filename)

        prs = Presentation()
        prs.slide_width  = Inches(self.SLIDE_W)
        prs.slide_height = Inches(self.SLIDE_H)

        chapters = book_info.get('chapters', [])
        self._slide_num = 0  # 슬라이드 번호 추적

        # ── 슬라이드 번호 사전 계산 (목차에 표시할 각 챕터 시작 슬라이드) ──
        slide_idx = 1  # 표지 = 슬라이드 1
        # 목차 슬라이드 수 계산
        half = (len(chapters) + 1) // 2
        toc_count = sum(1 for chunk in [chapters[:half], chapters[half:]] if chunk)
        slide_idx += toc_count

        # 프롤로그
        prologue_text = ebook_data.get('prologue', '')
        if prologue_text and prologue_text.strip():
            slide_idx += 1

        # 각 챕터 시작 슬라이드 번호 계산
        chapter_slide_nums = {}
        for i, ch_data in enumerate(ebook_data.get('chapters_content', [])):
            ch_num = chapters[i].get('chapter_num', i + 1) if i < len(chapters) else i + 1
            chapter_slide_nums[ch_num] = slide_idx + 1  # 1-based 표시 번호
            content = ch_data.get('content', '')
            sections = self._parse_sections(content)
            slide_idx += 1 + len(sections)  # 챕터 인트로 + 섹션 슬라이드

        # 1. 표지
        self._slide_cover(prs, title, subtitle, len(chapters))

        # 2. 목차 (슬라이드 번호 포함)
        self._slide_toc(prs, chapters, chapter_slide_nums)

        # 3. 프롤로그
        prologue = ebook_data.get('prologue', '')
        if prologue and prologue.strip():
            self._slide_text_page(prs, '프롤로그', prologue, '#1a1a2e')

        # 4. 챕터별 슬라이드
        for i, ch_data in enumerate(ebook_data.get('chapters_content', [])):
            chapter  = ch_data.get('chapter', {})
            content  = ch_data.get('content', '')
            phase    = chapter.get('phase', '')
            ch_title = chapter.get('title', '')
            col_hex  = PHASE_COLORS.get(phase, '#6c5ce7')
            ch_label = f"CH.{i+1:02d}  {ch_title[:45]}"

            self._slide_chapter_intro(prs, chapter, i + 1, col_hex)

            for sec in self._parse_sections(content):
                sec_title = sec['title'] or ch_title
                sec_lines = [l for l in sec['lines'] if l.strip()]

                if sec['type'] == 'highlight':
                    bullets = self._extract_bullets(sec_lines, max_bullets=4, max_chars=90)
                    self._slide_highlight(prs, sec_title, bullets, col_hex)
                elif sec['type'] == 'summary':
                    points = self._extract_bullets(sec_lines, max_bullets=6, max_chars=90)
                    self._slide_summary(prs, sec_title, points, col_hex)
                elif sec['type'] == 'checklist':
                    items = self._extract_bullets(sec_lines, max_bullets=7, max_chars=90)
                    self._slide_checklist(prs, sec_title, items, col_hex)
                else:
                    bullets = self._extract_bullets(sec_lines, max_bullets=5, max_chars=95)
                    self._slide_section(prs, sec_title, bullets, ch_label, col_hex)

        # 5. 에필로그
        epilogue = ebook_data.get('epilogue', '')
        if epilogue and epilogue.strip():
            self._slide_text_page(prs, '에필로그', epilogue, '#0a0a1a')

        # 슬라이드 번호 추가 (표지 제외)
        for idx, slide in enumerate(prs.slides):
            if idx == 0:
                continue
            self._add_slide_number(slide, idx + 1)

        prs.save(filepath)
        return filepath, filename
